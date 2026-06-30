from __future__ import annotations

import json
import mimetypes
import re
import subprocess
from datetime import datetime, timezone
from pathlib import Path
from uuid import uuid4

import httpx
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..models import AgentSession, Project
from .deck_design import TEMPLATE_ID, TEMPLATE_VERSION, apply_consulting_template, normalize_architecture_diagram


PROJECT_ROOT = Path(__file__).resolve().parents[3]
RUNTIME_DIR = PROJECT_ROOT / "scripts" / "ppt-runtime"
ENGINE_SCRIPT = RUNTIME_DIR / "generate-deck.mjs"
GENERATED_DIR = PROJECT_ROOT / "backend" / "generated"
RENDERER_VERSION = "consulting-default-renderer-9-medical-academic"


def export_pptx(project: Project, session: AgentSession) -> Path:
    return ensure_export_artifact(project, session)["pptx_path"]


def ensure_export_artifact(project: Project, session: AgentSession) -> dict:
    cached = _find_cached_artifact(project, session)
    if cached:
        return cached

    artifacts = session.artifacts or {}
    outline_artifact = apply_consulting_template(artifacts.get("outline", {}))
    outline = outline_artifact.get("slides", [])
    if not outline:
        raise RuntimeError("当前会话没有可导出的大纲")

    export_id = f"{project.id}-{uuid4().hex[:8]}"
    work_dir = GENERATED_DIR / export_id
    asset_dir = work_dir / "assets"
    qa_dir = work_dir / "qa"
    asset_dir.mkdir(parents=True, exist_ok=True)
    qa_dir.mkdir(parents=True, exist_ok=True)
    output = GENERATED_DIR / f"{session.id}.pptx"

    citations = artifacts.get("verify", {}).get("citations", [])
    images = artifacts.get("images", {}).get("items", [])
    local_case_images = _local_case_images(artifacts)
    slides = []
    for index, slide in enumerate(outline):
        citation_indices = slide.get("citation_indices") or []
        slide_citations = [
            citations[citation_index - 1]
            for citation_index in citation_indices
            if isinstance(citation_index, int) and 0 < citation_index <= len(citations)
        ]
        if not slide_citations and citations:
            citation_count = min(3 if slide.get("layout") in {"evidence", "data"} else 1, len(citations))
            start = index % len(citations)
            slide_citations = [
                citations[(start + offset) % len(citations)]
                for offset in range(citation_count)
            ]
        layout = "cover" if index == 0 else slide.get("layout") or slide.get("type") or "content"
        uses_image = layout in {"cover", "image_split", "case", "content"}
        assignment = slide.get("image_assignment") or {}
        image = None
        image_path = None
        if assignment.get("mode") == "none":
            uses_image = False
        elif assignment.get("mode") == "upload":
            image = {
                "description": assignment.get("filename") or slide.get("title"),
                "author": assignment.get("author") or "用户上传",
                "source_url": "",
            }
            local_path = Path(assignment.get("path", ""))
            image_path = local_path if local_path.exists() else None
        elif assignment.get("mode") == "search":
            image = {
                "url": assignment.get("url"),
                "description": assignment.get("description") or slide.get("image_query") or slide.get("title"),
                "author": assignment.get("author"),
                "source_url": assignment.get("source_url"),
            }
            image_path = _download_image(image, asset_dir, index) if image.get("url") else None
        else:
            local_image = _local_case_image_for_slide(slide, local_case_images)
            if local_image and uses_image:
                image = local_image
                local_path = Path(local_image.get("path", ""))
                image_path = local_path if local_path.exists() else None
            elif images and uses_image:
                image = images[index % len(images)]
                image_path = _download_image(image, asset_dir, index)
        if image_path and (image or {}).get("path"):
            image = {
                **(image or {}),
                "description": (image or {}).get("description") or Path(image_path).name,
                "author": (image or {}).get("author") or "用户上传病例材料",
                "source_url": "",
            }
        slides.append({
            **slide,
            "number": slide.get("number", index + 1),
            "citations": slide_citations,
            "image_path": str(image_path) if image_path else None,
            "image_alt": (image or {}).get("description") or slide.get("image_query") or slide.get("title"),
            "image_content_type": assignment.get("content_type") or mimetypes.guess_type(str(image_path))[0] if image_path else None,
            "image_source": (image or {}).get("source_url"),
            "image_author": (image or {}).get("author"),
        })

    payload = {
        "project": {
            "id": project.id,
            "title": project.title,
            "topic": project.topic,
            "slide_count": project.slide_count,
            "speaker_notes_enabled": project.speaker_notes_enabled,
        },
        "slides": slides,
        "citations": citations,
    }
    input_path = work_dir / "deck-input.json"
    input_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

    result = subprocess.run(
        ["node", str(ENGINE_SCRIPT), str(input_path), str(output), str(qa_dir)],
        cwd=RUNTIME_DIR,
        capture_output=True,
        text=True,
        timeout=180,
        check=False,
    )
    if result.returncode != 0:
        details = (result.stderr or result.stdout or "").strip()
        if "@oai/artifact-tool" not in details and "ERR_MODULE_NOT_FOUND" not in details:
            raise RuntimeError(f"专业 PPT 引擎导出失败：{details[:1200]}")
        _export_with_python_fallback(payload, output, qa_dir)
    else:
        _render_fallback_qa(payload, qa_dir, write_layout=False)
    if not output.exists() or output.stat().st_size < 10_000:
        raise RuntimeError("专业 PPT 引擎没有生成有效的 PPTX 文件")

    _validate_qa(qa_dir, len(slides))
    manifest = _artifact_manifest(project, session, output, work_dir, qa_dir, len(slides))
    (work_dir / "export-manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return _hydrate_manifest(manifest)


def _find_cached_artifact(project: Project, session: AgentSession) -> dict | None:
    if not GENERATED_DIR.exists():
        return None
    manifests = sorted(
        GENERATED_DIR.glob(f"{project.id}-*/export-manifest.json"),
        key=lambda path: path.stat().st_mtime,
        reverse=True,
    )
    for path in manifests:
        try:
            manifest = json.loads(path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            continue
        if manifest.get("session_id") != session.id:
            continue
        if manifest.get("export_revision") != (session.artifacts or {}).get("_export_revision"):
            continue
        if manifest.get("renderer_version") != RENDERER_VERSION:
            continue
        if manifest.get("design_template") != TEMPLATE_ID or manifest.get("template_version") != TEMPLATE_VERSION:
            continue
        hydrated = _hydrate_manifest(manifest)
        pptx_path = hydrated["pptx_path"]
        qa_dir = hydrated["qa_dir"]
        if pptx_path.exists() and qa_dir.exists():
            try:
                _validate_qa(qa_dir, int(manifest.get("slide_count", 0)))
            except RuntimeError:
                continue
            return hydrated
    return None


def _artifact_manifest(project: Project, session: AgentSession, output: Path, work_dir: Path, qa_dir: Path, slide_count: int) -> dict:
    slides = [
        {
            "number": index,
            "thumbnail_path": str((qa_dir / f"slide-{index:02d}.png").resolve()),
            "layout_path": str((qa_dir / f"slide-{index:02d}.layout.json").resolve()),
        }
        for index in range(1, slide_count + 1)
    ]
    return {
        "project_id": project.id,
        "session_id": session.id,
        "export_revision": (session.artifacts or {}).get("_export_revision"),
        "design_template": TEMPLATE_ID,
        "template_version": TEMPLATE_VERSION,
        "renderer_version": RENDERER_VERSION,
        "project_title": project.title,
        "slide_count": slide_count,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "file_status": "ready",
        "pptx_size_bytes": output.stat().st_size if output.exists() else 0,
        "pptx_path": str(output.resolve()),
        "work_dir": str(work_dir.resolve()),
        "qa_dir": str(qa_dir.resolve()),
        "montage_path": str((qa_dir / "montage.webp").resolve()),
        "warnings": build_export_warnings(session.artifacts or {}),
        "slides": slides,
    }


def _hydrate_manifest(manifest: dict) -> dict:
    hydrated = dict(manifest)
    hydrated["pptx_path"] = Path(hydrated["pptx_path"])
    hydrated["work_dir"] = Path(hydrated["work_dir"])
    hydrated["qa_dir"] = Path(hydrated["qa_dir"])
    hydrated["montage_path"] = Path(hydrated["montage_path"])
    hydrated["slides"] = [
        {
            **slide,
            "thumbnail_path": Path(slide["thumbnail_path"]),
            "layout_path": Path(slide["layout_path"]),
        }
        for slide in hydrated.get("slides", [])
    ]
    return hydrated


def _case_tokens(*values: str) -> set[str]:
    text = " ".join(str(value or "") for value in values)
    return {f"case-{int(match):02d}" for match in re.findall(r"case[-_\s]?0*(\d{1,2})", text, flags=re.IGNORECASE)}


def _local_case_images(artifacts: dict) -> list[dict]:
    items = artifacts.get("parse", {}).get("items", [])
    images = []
    for item in items:
        content_type = item.get("content_type") or ""
        path = item.get("path") or ""
        if not content_type.startswith("image/") or not path:
            continue
        descriptor = " ".join([
            item.get("filename", ""),
            item.get("description", ""),
            path,
        ])
        tokens = _case_tokens(descriptor)
        images.append({
            "path": path,
            "filename": item.get("filename", ""),
            "description": item.get("description") or item.get("filename", ""),
            "content_type": content_type,
            "author": "用户上传病例材料",
            "case_tokens": tokens,
        })
    return images


def _local_case_image_for_slide(slide: dict, local_images: list[dict]) -> dict | None:
    if not local_images:
        return None
    card = slide.get("case_card") or {}
    tokens = _case_tokens(
        slide.get("title", ""),
        slide.get("key_message", ""),
        " ".join(slide.get("bullets") or []),
        card.get("case_id", ""),
        card.get("source_filename", ""),
    )
    if tokens:
        for image in local_images:
            if tokens & image.get("case_tokens", set()):
                return image
    return None


def public_export_manifest(project: Project, session: AgentSession) -> dict:
    artifact = ensure_export_artifact(project, session)
    return {
        "session_id": artifact["session_id"],
        "project_id": artifact["project_id"],
        "project_title": artifact["project_title"],
        "slide_count": artifact["slide_count"],
        "created_at": artifact["created_at"],
        "file_status": artifact.get("file_status", "ready"),
        "pptx_size_bytes": artifact.get("pptx_size_bytes", 0),
        "design_template": artifact.get("design_template", TEMPLATE_ID),
        "template_version": artifact.get("template_version", TEMPLATE_VERSION),
        "renderer_version": artifact.get("renderer_version", RENDERER_VERSION),
        "pptx_url": f"/api/v1/sessions/{session.id}/export",
        "montage_url": f"/api/v1/sessions/{session.id}/export/montage",
        "warnings": artifact.get("warnings", []),
        "slides": [
            {
                "number": slide["number"],
                "thumbnail_url": f"/api/v1/sessions/{session.id}/export/thumbnails/{slide['number']}",
            }
            for slide in artifact.get("slides", [])
        ],
    }


def build_export_warnings(artifacts: dict) -> list[dict]:
    warnings = []
    slides = artifacts.get("outline", {}).get("slides", [])
    citations = artifacts.get("verify", {}).get("citations", [])
    images = artifacts.get("images", {}).get("items", [])
    parse = artifacts.get("parse", {})

    if not citations:
        warnings.append({
            "code": "missing_citations",
            "severity": "warning",
            "message": "当前导出没有可验证引用，建议补充文献或网页来源。",
            "slide_number": None,
        })

    if parse.get("ocr_required"):
        items = "、".join(parse.get("ocr_items") or parse.get("ocr_slides") or [])
        warnings.append({
            "code": "ocr_not_executed",
            "severity": "warning",
            "message": f"存在需要 OCR/视觉理解的资料{f'：{items}' if items else ''}，当前尚未执行 OCR。",
            "slide_number": None,
        })

    for index, slide in enumerate(slides):
        layout = "cover" if index == 0 else slide.get("layout") or slide.get("type") or "content"
        assignment = slide.get("image_assignment") or {}
        needs_image = layout in {"cover", "image_split", "case", "content"}
        has_auto_image = bool(images)
        has_selected_image = assignment.get("mode") in {"search", "upload"}
        if assignment.get("mode") == "none":
            has_selected_image = False
            has_auto_image = False
        if needs_image and not has_auto_image and not has_selected_image:
            warnings.append({
                "code": "missing_image",
                "severity": "warning",
                "message": f"第 {index + 1} 页使用图文版式但没有可用图片。",
                "slide_number": index + 1,
            })
        if assignment.get("mode") == "upload" and assignment.get("path") and not Path(assignment["path"]).exists():
            warnings.append({
                "code": "image_file_missing",
                "severity": "error",
                "message": f"第 {index + 1} 页上传图片文件不存在，导出会降级为空占位。",
                "slide_number": index + 1,
            })
    return warnings


def export_thumbnail_path(project: Project, session: AgentSession, slide_number: int) -> Path:
    artifact = ensure_export_artifact(project, session)
    for slide in artifact.get("slides", []):
        if slide["number"] == slide_number and slide["thumbnail_path"].exists():
            return slide["thumbnail_path"]
    raise RuntimeError("缩略图不存在")


def export_montage_path(project: Project, session: AgentSession) -> Path:
    artifact = ensure_export_artifact(project, session)
    if artifact["montage_path"].exists():
        return artifact["montage_path"]
    raise RuntimeError("总览图不存在")


def _download_image(image: dict, asset_dir: Path, index: int) -> Path | None:
    url = image.get("url")
    if not url:
        return None
    try:
        with httpx.Client(timeout=25, follow_redirects=True) as client:
            response = client.get(url)
            response.raise_for_status()
        content_type = response.headers.get("content-type", "image/jpeg").split(";")[0]
        extension = mimetypes.guess_extension(content_type) or ".jpg"
        target = asset_dir / f"slide-{index + 1:02d}{extension}"
        target.write_bytes(response.content)
        return target
    except Exception:
        return None


def _validate_qa(qa_dir: Path, slide_count: int) -> None:
    pngs = sorted(qa_dir.glob("slide-*.png"))
    layouts = sorted(qa_dir.glob("slide-*.layout.json"))
    montage = qa_dir / "montage.webp"
    if len(pngs) != slide_count or len(layouts) != slide_count or not montage.exists():
        raise RuntimeError("PPT 渲染 QA 输出不完整")
    for png in pngs:
        if png.stat().st_size < 2_000:
            raise RuntimeError(f"幻灯片渲染结果异常：{png.name}")


def _export_with_python_fallback(payload: dict, output: Path, qa_dir: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = payload.get("slides", [])
    for index, item in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        _pptx_background(slide, "FFFFFF")
        layout = item.get("layout") or item.get("type") or "content"
        if index == 0 and layout == "cover":
            _add_textbox(slide, "PPTKILLER / FALLBACK EXPORT", 0.7, 0.55, 4.6, 0.3, 12, "5442D6", True)
            _add_textbox(slide, item.get("title") or payload["project"]["title"], 0.7, 1.75, 7.8, 1.45, 34, "111936", True)
            _add_textbox(slide, payload["project"].get("topic") or item.get("key_message") or "", 0.72, 3.32, 7.6, 0.7, 18, "687086")
            _add_textbox(slide, f"{payload['project'].get('slide_count', len(slides))} 页 · 可编辑 PPTX", 0.72, 5.72, 4.5, 0.3, 11, "687086")
        elif layout == "two_column":
            _add_textbox(slide, item.get("title", f"第 {index + 1} 页"), 0.7, 0.75, 9.6, 0.58, 25, "111936", True)
            _add_two_column_content(slide, item)
        elif layout == "process":
            _add_textbox(slide, item.get("title", f"第 {index + 1} 页"), 0.7, 0.75, 9.6, 0.58, 25, "111936", True)
            _add_process_content(slide, item)
        elif layout == "architecture":
            _add_textbox(slide, item.get("title", "架构图"), 0.7, 0.75, 8.8, 0.58, 25, "111936", True)
            _add_textbox(slide, item.get("key_message", "用结构图说明模块、关系与反馈闭环"), 0.72, 1.36, 9.0, 0.38, 12, "687086")
            _add_architecture_diagram(slide, item)
        elif layout == "data":
            _add_textbox(slide, item.get("title", "数据页"), 0.7, 0.75, 8.8, 0.58, 25, "111936", True)
            _add_textbox(slide, item.get("key_message", "数据支撑核心判断"), 0.72, 1.36, 9.0, 0.38, 12, "687086")
            _add_data_table(slide, item)
        else:
            _add_textbox(slide, item.get("title", f"第 {index + 1} 页"), 0.7, 0.75, 9.6, 0.7, 26, "111936", True)
            _add_textbox(slide, item.get("key_message", ""), 0.72, 1.47, 9.0, 0.45, 14, "5442D6", True)
            bullets = "\n".join(f"• {bullet}" for bullet in (item.get("bullets") or [])[:5])
            _add_textbox(slide, bullets, 0.9, 2.1, 7.4, 3.2, 16, "111936")
            if item.get("image_path"):
                try:
                    slide.shapes.add_picture(item["image_path"], Inches(8.7), Inches(1.8), width=Inches(3.7), height=Inches(3.2))
                except Exception:
                    _add_placeholder(slide, item.get("image_alt") or "图片占位")
        footer = _footer_text(item, index, len(slides))
        _add_textbox(slide, footer, 0.7, 7.05, 10.3, 0.22, 7, "687086")
        _add_textbox(slide, f"{index + 1} / {len(slides)}", 11.2, 7.05, 1.2, 0.22, 8, "687086")
        if payload.get("project", {}).get("speaker_notes_enabled") and item.get("speaker_notes"):
            slide.notes_slide.notes_text_frame.text = str(item["speaker_notes"])
    prs.save(output)
    _render_fallback_qa(payload, qa_dir)


def _pptx_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def _add_textbox(slide, text: str, left: float, top: float, width: float, height: float, size: int, color: str, bold: bool = False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def _add_placeholder(slide, label: str) -> None:
    shape = slide.shapes.add_shape(1, Inches(8.7), Inches(1.8), Inches(3.7), Inches(3.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string("F1F2F8")
    shape.line.color.rgb = RGBColor.from_string("D9DDE8")
    _add_textbox(slide, label, 9.05, 3.05, 3.0, 0.4, 12, "687086")


def _add_two_column_content(slide, item: dict) -> None:
    columns = [
        (0.85, item.get("left_title") or "现状与机会", item.get("left_bullets") or (item.get("bullets") or [])[:3], "F8FAFD", "5442D6"),
        (6.9, item.get("right_title") or "挑战与应对", item.get("right_bullets") or (item.get("bullets") or [])[3:6], "FFFFFF", "14245A"),
    ]
    for left, title, bullets, fill, title_color in columns:
        panel = slide.shapes.add_shape(1, Inches(left), Inches(2.05), Inches(5.35), Inches(3.8))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor.from_string(fill)
        panel.line.color.rgb = RGBColor.from_string("D9DDE8")
        _add_textbox(slide, title, left + 0.3, 2.38, 4.7, 0.34, 15, title_color, True)
        bullet_text = "\n\n".join(f"• {bullet}" for bullet in [str(item) for item in bullets if str(item).strip()][:4])
        _add_textbox(slide, bullet_text, left + 0.35, 3.05, 4.5, 2.1, 13, "111936")


def _add_process_content(slide, item: dict) -> None:
    steps = [str(step) for step in (item.get("process_steps") or item.get("bullets") or []) if str(step).strip()][:5]
    if not steps and item.get("key_message"):
        steps = [str(item["key_message"])]
    gap = 0.22
    width = (11.2 - gap * max(len(steps) - 1, 0)) / max(len(steps), 1)
    for index, step in enumerate(steps):
        left = 0.82 + index * (width + gap)
        _add_textbox(slide, str(index + 1).zfill(2), left, 2.25, width, 0.45, 23, "5442D6", True)
        rule = slide.shapes.add_shape(1, Inches(left), Inches(2.95), Inches(width), Inches(0.03))
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor.from_string("5442D6" if index == len(steps) - 1 else "D9DDE8")
        rule.line.color.rgb = rule.fill.fore_color.rgb
        panel = slide.shapes.add_shape(1, Inches(left), Inches(3.28), Inches(width), Inches(1.7))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor.from_string("F8FAFD")
        panel.line.color.rgb = RGBColor.from_string("D9DDE8")
        _add_textbox(slide, step, left + 0.14, 3.55, max(width - 0.28, 0.5), 0.9, 12, "111936", True)


def _add_architecture_diagram(slide, item: dict) -> None:
    normalized = normalize_architecture_diagram(item)
    modules = (normalized.get("diagram_modules") or [])[:5]
    if not modules:
        return
    _add_textbox(slide, normalized.get("diagram_title") or "主链路架构", 0.75, 2.0, 4.5, 0.32, 14, "5442D6", True)
    gap = 0.18
    module_width = (10.9 - gap * (len(modules) - 1)) / len(modules)
    top = 2.65
    for module_index, module in enumerate(modules):
        left = 0.78 + module_index * (module_width + gap)
        panel = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(module_width), Inches(2.72))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor.from_string("F8FAFD")
        panel.line.color.rgb = RGBColor.from_string("D9DDE8")
        _add_textbox(slide, str(module_index + 1).zfill(2), left + 0.16, top + 0.16, 0.34, 0.18, 8, "5442D6", True)
        _add_textbox(slide, module.get("label", ""), left + 0.16, top + 0.46, module_width - 0.32, 0.34, 11, "111936", True)
        for child_index, child in enumerate((module.get("children") or [])[:2]):
            child_top = top + 1.02 + child_index * 0.62
            box = slide.shapes.add_shape(1, Inches(left + 0.14), Inches(child_top), Inches(module_width - 0.28), Inches(0.45))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
            box.line.color.rgb = RGBColor.from_string("D9DDE8")
            _add_textbox(slide, child.get("label", ""), left + 0.24, child_top + 0.1, module_width - 0.48, 0.15, 7, "111936", True)
            _add_textbox(slide, child.get("detail", ""), left + 0.24, child_top + 0.25, module_width - 0.48, 0.12, 5, "687086")
        if module_index < len(modules) - 1:
            line = slide.shapes.add_shape(1, Inches(left + module_width + 0.04), Inches(top + 1.34), Inches(gap - 0.06), Inches(0.02))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor.from_string("5442D6")
            line.line.color.rgb = RGBColor.from_string("5442D6")
            _add_textbox(slide, ">", left + module_width + gap - 0.12, top + 1.22, 0.16, 0.14, 8, "5442D6", True)


def _add_data_table(slide, item: dict) -> None:
    headers = (item.get("data_table") or {}).get("headers", [])[:5]
    rows = (item.get("data_table") or {}).get("rows", [])[:6]
    if not headers:
        kpis = item.get("kpis") or []
        for index, kpi in enumerate(kpis[:4]):
            _add_textbox(slide, kpi.get("label", "指标"), 0.9 + index * 2.8, 2.2, 2.3, 0.35, 11, "687086")
            _add_textbox(slide, kpi.get("value", ""), 0.9 + index * 2.8, 2.65, 2.3, 0.55, 25, "5442D6", True)
        return
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.75), Inches(2.05), Inches(11.2), Inches(3.8)).table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string("18256F")
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string("FFFFFF")
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
    for row_index, row in enumerate(rows, start=1):
        for col_index, _ in enumerate(headers):
            table.cell(row_index, col_index).text = str(row[col_index] if col_index < len(row) else "")
    source = item.get("data_source") or {}
    _add_textbox(slide, f"数据来源：{source.get('filename', '')} / {source.get('sheet', '')}", 0.75, 6.15, 6.5, 0.25, 9, "687086")


def _footer_text(item: dict, index: int, total: int) -> str:
    citations = (item.get("citations") or [])[:2]
    if citations:
        return " · ".join(citation.get("doi") or citation.get("url") or citation.get("title", "") for citation in citations)
    if item.get("image_author"):
        return f"图片：{item.get('image_author')}"
    return "PPTKiller · fallback editable export"


def _render_fallback_qa(payload: dict, qa_dir: Path, write_layout: bool = True) -> None:
    qa_dir.mkdir(parents=True, exist_ok=True)
    thumbnails = []
    for index, item in enumerate(payload.get("slides", []), start=1):
        image = Image.new("RGB", (1280, 720), "white")
        draw = ImageDraw.Draw(image)
        font_title = _fallback_font(38)
        font_body = _fallback_font(22)
        draw.rectangle((0, 0, 18, 720), fill="#5442D6")
        draw.text((74, 78), str(index).zfill(2), fill="#5442D6", font=_fallback_font(18))
        draw.text((74, 128), item.get("title", f"第 {index} 页"), fill="#111936", font=font_title)
        layout = item.get("layout") or item.get("type") or "content"
        if layout == "two_column":
            _draw_preview_two_column(draw, item)
        elif layout == "process":
            _draw_preview_process(draw, item)
        elif layout == "architecture":
            normalized = normalize_architecture_diagram(item)
            modules = (normalized.get("diagram_modules") or [])[:5]
            draw.text((96, 230), normalized.get("diagram_title", "主链路架构"), fill="#5442D6", font=font_body)
            gap = 18
            module_width = (1030 - gap * max(len(modules) - 1, 0)) // max(len(modules), 1)
            for module_index, module in enumerate(modules):
                x = 96 + module_index * (module_width + gap)
                draw.rectangle((x, 294, x + module_width, 560), outline="#D9DDE8", fill="#F8FAFD")
                draw.text((x + 18, 316), str(module_index + 1).zfill(2), fill="#5442D6", font=_fallback_font(15))
                draw.text((x + 18, 350), module.get("label", ""), fill="#111936", font=_fallback_font(18))
                for child_index, child in enumerate((module.get("children") or [])[:2]):
                    y_node = 414 + child_index * 58
                    draw.rectangle((x + 16, y_node, x + module_width - 16, y_node + 42), outline="#D9DDE8", fill="#FFFFFF")
                    draw.text((x + 28, y_node + 12), child.get("label", ""), fill="#111936", font=_fallback_font(13))
                if module_index < len(modules) - 1:
                    y_mid = 426
                    draw.line((x + module_width + 4, y_mid, x + module_width + gap - 4, y_mid), fill="#5442D6", width=3)
                    draw.text((x + module_width + gap - 14, y_mid - 14), ">", fill="#5442D6", font=_fallback_font(16))
        elif layout == "data":
            draw.text((96, 230), item.get("key_message", "数据页"), fill="#5442D6", font=font_body)
            headers = (item.get("data_table") or {}).get("headers", [])[:5]
            rows = (item.get("data_table") or {}).get("rows", [])[:6]
            draw.text((96, 300), " | ".join(map(str, headers)), fill="#18256F", font=font_body)
            for row_index, row in enumerate(rows):
                draw.text((96, 342 + row_index * 38), " | ".join(map(str, row[:5])), fill="#111936", font=_fallback_font(17))
        else:
            _draw_preview_bullets(draw, item)
        draw.text((74, 674), f"{index} / {len(payload.get('slides', []))}", fill="#687086", font=_fallback_font(14))
        png = qa_dir / f"slide-{index:02d}.png"
        image.save(png)
        thumbnails.append(image.resize((320, 180)))
        if write_layout:
            (qa_dir / f"slide-{index:02d}.layout.json").write_text(
                json.dumps({"fallback": True, "slide": index, "title": item.get("title", "")}, ensure_ascii=False),
                encoding="utf-8",
            )
    if thumbnails:
        columns = min(4, len(thumbnails))
        rows = (len(thumbnails) + columns - 1) // columns
        montage = Image.new("RGB", (columns * 320, rows * 180), "#F1F2F8")
        for index, thumb in enumerate(thumbnails):
            montage.paste(thumb, ((index % columns) * 320, (index // columns) * 180))
        montage.save(qa_dir / "montage.webp")


def _draw_preview_bullets(draw: ImageDraw.ImageDraw, item: dict) -> None:
    bullets = [str(item) for item in (item.get("bullets") or []) if str(item).strip()][:5]
    if not bullets and item.get("key_message"):
        bullets = [str(item["key_message"])]
    for index, bullet in enumerate(bullets):
        top = 238 + index * 64
        draw.rounded_rectangle((96, top, 850, top + 46), radius=8, outline="#E3E7F0", fill="#F8FAFD")
        draw.text((120, top + 12), f"• {bullet}", fill="#111936", font=_fallback_font(20))


def _draw_preview_two_column(draw: ImageDraw.ImageDraw, item: dict) -> None:
    columns = [
        (96, item.get("left_title") or "现状与机会", item.get("left_bullets") or (item.get("bullets") or [])[:3], "#F8FAFD"),
        (670, item.get("right_title") or "挑战与应对", item.get("right_bullets") or (item.get("bullets") or [])[3:6], "#FFFFFF"),
    ]
    for left, title, bullets, fill in columns:
        draw.rounded_rectangle((left, 230, left + 510, 585), radius=10, outline="#E3E7F0", fill=fill)
        draw.text((left + 28, 262), str(title), fill="#5442D6" if left == 96 else "#14245A", font=_fallback_font(24))
        for index, bullet in enumerate([str(item) for item in bullets if str(item).strip()][:4]):
            y = 324 + index * 52
            draw.text((left + 34, y), f"• {bullet}", fill="#111936", font=_fallback_font(18))


def _draw_preview_process(draw: ImageDraw.ImageDraw, item: dict) -> None:
    steps = [str(item) for item in (item.get("process_steps") or item.get("bullets") or []) if str(item).strip()][:5]
    if not steps and item.get("key_message"):
        steps = [str(item["key_message"])]
    gap = 22
    count = max(len(steps), 1)
    width = int((1088 - gap * (count - 1)) / count)
    for index, step in enumerate(steps):
        left = 96 + index * (width + gap)
        draw.text((left, 250), str(index + 1).zfill(2), fill="#5442D6", font=_fallback_font(34))
        draw.line((left, 304, left + width, 304), fill="#5442D6" if index == count - 1 else "#E3E7F0", width=4)
        draw.rounded_rectangle((left, 334, left + width, 500), radius=8, outline="#E3E7F0", fill="#F8FAFD")
        _draw_wrapped_text(draw, step, left + 18, 362, width - 36, _fallback_font(18), "#111936", line_gap=28, max_lines=4)


def _draw_wrapped_text(draw: ImageDraw.ImageDraw, text: str, x: int, y: int, width: int, font, fill: str, line_gap: int = 24, max_lines: int = 3) -> None:
    line = ""
    lines = []
    for char in str(text):
        candidate = line + char
        if draw.textlength(candidate, font=font) <= width:
            line = candidate
        else:
            if line:
                lines.append(line)
            line = char
        if len(lines) >= max_lines:
            break
    if line and len(lines) < max_lines:
        lines.append(line)
    for index, chunk in enumerate(lines[:max_lines]):
        draw.text((x, y + index * line_gap), chunk, fill=fill, font=font)


def _fallback_font(size: int):
    for candidate in [
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Songti.ttc",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]:
        try:
            return ImageFont.truetype(candidate, size)
        except Exception:
            continue
    return ImageFont.load_default()
