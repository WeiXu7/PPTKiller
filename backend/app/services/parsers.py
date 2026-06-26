from __future__ import annotations

import csv
import io
from pathlib import Path
from typing import Any

from PIL import Image


MAX_TEXT_CHARS = 60_000
MAX_ROWS_PER_SHEET = 200
MAX_COLS_PER_SHEET = 40


class AssetParser:
    def parse(self, path: str, filename: str, content_type: str) -> dict[str, Any]:
        source = Path(path)
        suffix = Path(filename).suffix.lower()
        base = {
            "filename": filename,
            "content_type": content_type,
            "parser": "metadata",
            "sections": [],
            "text": "",
            "tables": [],
            "images": [],
            "ocr_required": False,
            "warnings": [],
        }
        try:
            if suffix == ".pdf" or content_type == "application/pdf":
                return {**base, **self._parse_pdf(source)}
            if suffix == ".docx":
                return {**base, **self._parse_docx(source)}
            if suffix == ".pptx":
                return {**base, **self._parse_pptx(source)}
            if suffix in {".xlsx", ".xlsm"}:
                return {**base, **self._parse_xlsx(source)}
            if suffix == ".xls":
                return {**base, **self._parse_xls(source)}
            if suffix in {".csv", ".tsv"}:
                return {**base, **self._parse_delimited(source, "\t" if suffix == ".tsv" else ",")}
            if suffix in {".txt", ".md", ".json", ".html", ".htm"} or content_type.startswith("text/"):
                text = source.read_text(encoding="utf-8", errors="ignore")[:MAX_TEXT_CHARS]
                return {
                    **base,
                    "parser": "text",
                    "text": text,
                    "sections": [{"kind": "text", "label": filename, "text": text}],
                    "stats": {"characters": len(text)},
                }
            if content_type.startswith("image/") or suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".tiff", ".bmp"}:
                return {**base, **self._parse_image(source)}
            return {**base, "warnings": ["暂不支持该文件格式的正文解析，仅保留文件元数据。"]}
        except Exception as exc:
            return {
                **base,
                "parser_error": f"{type(exc).__name__}: {str(exc)[:300]}",
                "warnings": ["文件解析失败，原始文件仍被保留。"],
            }

    def _parse_pdf(self, path: Path) -> dict:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        sections = []
        text_parts = []
        image_pages = 0
        scanned_pages = []
        for index, page in enumerate(reader.pages):
            text = (page.extract_text() or "").strip()
            image_count = len(getattr(page, "images", []) or [])
            if image_count:
                image_pages += 1
            if len(text) < 20 and image_count:
                scanned_pages.append(index + 1)
            sections.append({
                "kind": "page",
                "page": index + 1,
                "label": f"第 {index + 1} 页",
                "text": text[:8000],
                "image_count": image_count,
            })
            if text:
                text_parts.append(f"[第 {index + 1} 页]\n{text}")
        text = "\n\n".join(text_parts)[:MAX_TEXT_CHARS]
        return {
            "parser": "pypdf",
            "text": text,
            "sections": sections,
            "ocr_required": bool(scanned_pages),
            "ocr_pages": scanned_pages,
            "stats": {
                "pages": len(reader.pages),
                "text_characters": len(text),
                "pages_with_images": image_pages,
                "scanned_pages": len(scanned_pages),
            },
            "warnings": [f"{len(scanned_pages)} 页缺少有效文本层，建议按需 OCR。"] if scanned_pages else [],
        }

    def _parse_docx(self, path: Path) -> dict:
        from docx import Document

        document = Document(str(path))
        sections = []
        text_parts = []
        for index, paragraph in enumerate(document.paragraphs):
            text = paragraph.text.strip()
            if not text:
                continue
            sections.append({
                "kind": "paragraph",
                "index": index + 1,
                "style": paragraph.style.name if paragraph.style else "",
                "label": paragraph.style.name if paragraph.style else f"段落 {index + 1}",
                "text": text,
            })
            text_parts.append(text)
        tables = []
        for table_index, table in enumerate(document.tables):
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            tables.append({"index": table_index + 1, "rows": rows[:MAX_ROWS_PER_SHEET]})
            text_parts.append(f"[表格 {table_index + 1}]\n" + "\n".join(" | ".join(row) for row in rows[:50]))
        image_count = len(document.inline_shapes)
        return {
            "parser": "python-docx",
            "text": "\n\n".join(text_parts)[:MAX_TEXT_CHARS],
            "sections": sections,
            "tables": tables,
            "images": [{"kind": "embedded", "count": image_count}] if image_count else [],
            "ocr_required": image_count > 0 and not text_parts,
            "stats": {
                "paragraphs": len(sections),
                "tables": len(tables),
                "embedded_images": image_count,
            },
            "warnings": ["文档仅包含图片，建议 OCR。"] if image_count > 0 and not text_parts else [],
        }

    def _parse_pptx(self, path: Path) -> dict:
        from pptx import Presentation

        presentation = Presentation(str(path))
        sections = []
        text_parts = []
        tables = []
        total_images = 0
        image_only_slides = []
        for slide_index, slide in enumerate(presentation.slides):
            slide_text = []
            slide_images = 0
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    value = shape.text.strip()
                    if value:
                        slide_text.append(value)
                if getattr(shape, "shape_type", None) == 13:
                    slide_images += 1
                if getattr(shape, "has_table", False):
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    tables.append({"slide": slide_index + 1, "rows": rows[:MAX_ROWS_PER_SHEET]})
            notes = ""
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass
            total_images += slide_images
            if not slide_text and slide_images:
                image_only_slides.append(slide_index + 1)
            text = "\n".join(slide_text)
            sections.append({
                "kind": "slide",
                "slide": slide_index + 1,
                "label": slide_text[0][:120] if slide_text else f"幻灯片 {slide_index + 1}",
                "text": text[:8000],
                "notes": notes[:4000],
                "image_count": slide_images,
            })
            text_parts.append(f"[幻灯片 {slide_index + 1}]\n{text}" + (f"\n[备注]\n{notes}" if notes else ""))
        return {
            "parser": "python-pptx",
            "text": "\n\n".join(text_parts)[:MAX_TEXT_CHARS],
            "sections": sections,
            "tables": tables,
            "images": [{"kind": "embedded", "count": total_images}] if total_images else [],
            "ocr_required": bool(image_only_slides),
            "ocr_slides": image_only_slides,
            "stats": {
                "slides": len(presentation.slides),
                "tables": len(tables),
                "embedded_images": total_images,
                "image_only_slides": len(image_only_slides),
            },
            "warnings": [f"{len(image_only_slides)} 页只有图片，建议按需 OCR 或视觉理解。"] if image_only_slides else [],
        }

    def _parse_xlsx(self, path: Path) -> dict:
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), read_only=True, data_only=False)
        return self._workbook_payload(
            [(sheet.title, sheet.iter_rows(values_only=True), sheet.max_row, sheet.max_column) for sheet in workbook.worksheets],
            "openpyxl",
        )

    def _parse_xls(self, path: Path) -> dict:
        import xlrd

        workbook = xlrd.open_workbook(str(path))
        sheets = []
        for sheet in workbook.sheets():
            rows = ([sheet.cell_value(row, col) for col in range(sheet.ncols)] for row in range(sheet.nrows))
            sheets.append((sheet.name, rows, sheet.nrows, sheet.ncols))
        return self._workbook_payload(sheets, "xlrd")

    def _workbook_payload(self, sheets, parser: str) -> dict:
        sections = []
        tables = []
        text_parts = []
        for name, rows_iter, row_count, col_count in sheets:
            rows = []
            for row_index, row in enumerate(rows_iter):
                if row_index >= MAX_ROWS_PER_SHEET:
                    break
                values = [self._cell_value(value) for value in list(row)[:MAX_COLS_PER_SHEET]]
                if any(value != "" for value in values):
                    rows.append(values)
            tables.append({"sheet": name, "rows": rows})
            preview = "\n".join(" | ".join(row) for row in rows[:60])
            sections.append({
                "kind": "sheet",
                "label": name,
                "sheet": name,
                "text": preview[:10000],
                "rows": row_count,
                "columns": col_count,
            })
            text_parts.append(f"[工作表：{name}]\n{preview}")
        return {
            "parser": parser,
            "text": "\n\n".join(text_parts)[:MAX_TEXT_CHARS],
            "sections": sections,
            "tables": tables,
            "stats": {"sheets": len(sections), "tables": len(tables)},
        }

    def _parse_delimited(self, path: Path, delimiter: str) -> dict:
        text = path.read_text(encoding="utf-8-sig", errors="ignore")
        rows = []
        for index, row in enumerate(csv.reader(io.StringIO(text), delimiter=delimiter)):
            if index >= MAX_ROWS_PER_SHEET:
                break
            rows.append(row[:MAX_COLS_PER_SHEET])
        preview = "\n".join(" | ".join(row) for row in rows)
        return {
            "parser": "csv",
            "text": preview[:MAX_TEXT_CHARS],
            "sections": [{"kind": "sheet", "label": path.name, "text": preview[:10000]}],
            "tables": [{"sheet": path.stem, "rows": rows}],
            "stats": {"rows": len(rows), "columns": max((len(row) for row in rows), default=0)},
        }

    def _parse_image(self, path: Path) -> dict:
        with Image.open(path) as image:
            return {
                "parser": "pillow",
                "images": [{
                    "kind": "source",
                    "width": image.width,
                    "height": image.height,
                    "format": image.format,
                    "mode": image.mode,
                }],
                "ocr_required": True,
                "stats": {"width": image.width, "height": image.height, "format": image.format},
                "warnings": ["图片没有附带文本描述，后续需要视觉理解或 OCR。"],
            }

    @staticmethod
    def _cell_value(value: Any) -> str:
        if value is None:
            return ""
        return str(value)[:1000]
