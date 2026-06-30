import asyncio
from pathlib import Path

from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

from backend.app.database import Base
from backend.app.models import Project, ProjectAsset, User
from backend.app.security import hash_password, verify_password
from backend.app.api import apply_slide_update, apply_slide_image_assignment, refresh_case_preparation
from backend.app.schemas import SlideUpdateRequest
from backend.app.services.agent_services import AgentServices
from backend.app.services.deck_design import apply_consulting_template, normalize_architecture_diagram
from backend.app.services.harness import AgentHarness
from backend.app.services.providers import Citation, _pubmed_doi, _year_from_pubdate, dedupe_citations
from backend.app.services.pptx_export import (
    build_export_warnings, export_pptx, public_export_manifest,
    _export_with_python_fallback, _local_case_image_for_slide, _local_case_images, _render_fallback_qa,
)


class FakeServices:
    async def parse_assets(self, project):
        return {"count": 0, "files": [], "text_context": "", "note": "无上传资料"}

    async def case_prepare(self, project, brief, parsed):
        return {
            "mode": "radiology_case",
            "readiness": "not_recommended",
            "readiness_label": "不建议直接生成",
            "case_count": 0,
            "asset_count": 0,
            "case_collection_checklist": [],
            "case_cards": [],
            "next_questions": ["请补充已脱敏病例材料。"],
            "note": "已生成病例准备清单",
        }

    async def research(self, project):
        return {
            "configured": True,
            "provider": "Fake",
            "found": 1,
            "peer_reviewed": 1,
            "web_results": [{"title": "Source", "url": "https://example.com"}],
            "citations": [{
                "title": "Verified Paper",
                "authors": ["A. Author"],
                "year": 2025,
                "venue": "Journal",
                "url": "https://doi.org/10.1000/test",
                "doi": "10.1000/test",
                "verified": True,
            }],
            "note": "真实服务测试替身",
        }

    async def verify(self, research):
        return {
            "verified": 1,
            "doi_checked": 1,
            "duplicates_removed": 0,
            "citations": research["citations"],
            "note": "已验证",
        }

    async def images(self, project):
        return {"configured": True, "provider": "Fake", "found": 1, "selected": 1, "items": [], "note": "已检索图片"}

    async def outline(self, project, brief, artifacts):
        slides = []
        for index in range(project.slide_count):
            slides.append({
                "number": index + 1,
                "title": project.title if index == 0 else f"第 {index + 1} 页",
                "type": "cover" if index == 0 else "content",
                "bullets": ["核心观点", "证据与建议"],
                "citation_indices": [1],
                "speaker_notes": "演讲稿",
            })
        return {"slides": slides, "target_count": project.slide_count, "generated_by": "fake", "note": "已生成"}

    async def slides(self, project, artifacts):
        return {"generated": len(artifacts["outline"]["slides"]), "slides": artifacts["outline"]["slides"], "note": "已准备"}

    async def notes(self, project, artifacts):
        return {"enabled": True, "coverage": project.slide_count, "items": [], "note": "已生成演讲稿"}


def test_password_roundtrip():
    encoded = hash_password("correct horse battery staple")
    assert verify_password("correct horse battery staple", encoded)
    assert not verify_password("wrong", encoded)


def test_harness_stops_for_outline_approval():
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(engine)
    db = sessionmaker(bind=engine)()
    user = User(email="demo@example.com", name="Demo", password_hash="x")
    db.add(user)
    db.flush()
    project = Project(owner_id=user.id, title="AI 教育趋势", slide_count=15)
    db.add(project)
    db.commit()
    session = asyncio.run(AgentHarness(db, FakeServices()).start(project, {"audience": "教育工作者"}))
    assert session.status == "waiting_approval"
    assert session.events[-1].step_key == "outline"


def test_harness_completes_after_two_approvals(tmp_path, monkeypatch):
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(engine)
    db = sessionmaker(bind=engine)()
    user = User(email="complete@example.com", name="Complete", password_hash="x")
    db.add(user)
    db.flush()
    project = Project(owner_id=user.id, title="完整闭环", topic="验证导出", slide_count=12)
    db.add(project)
    db.commit()
    harness = AgentHarness(db, FakeServices())
    session = asyncio.run(harness.start(project, {"audience": "管理者"}))
    session = asyncio.run(harness.approve(session, project, True, ""))
    assert session.status == "waiting_approval"
    assert session.current_step == 9
    session = asyncio.run(harness.approve(session, project, True, ""))
    assert session.status == "completed"
    assert project.status == "completed"
    monkeypatch.chdir(tmp_path)
    (tmp_path / "backend" / "generated").mkdir(parents=True)
    output = export_pptx(project, session)
    assert output.exists()
    assert output.stat().st_size > 0


def test_radiology_case_mode_stops_for_case_preparation():
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(engine)
    db = sessionmaker(bind=engine)()
    user = User(email="radiology@example.com", name="Radiology", password_hash="x")
    db.add(user)
    db.flush()
    project = Project(owner_id=user.id, title="肺结节影像诊断病例教学", topic="病例驱动教学", slide_count=10)
    db.add(project)
    db.commit()

    session = asyncio.run(AgentHarness(db, FakeServices()).start(project, {
        "audience": "影像科医生",
        "presentation_mode": "radiology_case",
    }))

    assert session.status == "waiting_approval"
    assert session.current_step == 2
    assert session.events[-1].step_key == "case_prepare"
    assert session.artifacts["case_prepare"]["readiness"] == "not_recommended"


def test_refresh_case_preparation_reparses_uploaded_case_material(tmp_path):
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(engine)
    db = sessionmaker(bind=engine)()
    user = User(email="case-refresh@example.com", name="Case Refresh", password_hash="x")
    db.add(user)
    db.flush()
    project = Project(owner_id=user.id, title="肺结节影像诊断病例教学", topic="病例驱动教学", slide_count=10)
    db.add(project)
    db.commit()
    session = asyncio.run(AgentHarness(db, FakeServices()).start(project, {
        "audience": "影像科医生",
        "presentation_mode": "radiology_case",
    }))

    case_file = tmp_path / "case-001.txt"
    case_file.write_text("65岁男性，CT 增强检查。最终诊断：肺腺癌，病理确认。材料已脱敏。", encoding="utf-8")
    db.add(ProjectAsset(
        project_id=project.id,
        filename=case_file.name,
        content_type="text/plain",
        path=str(case_file),
        description="典型病例",
        size_bytes=case_file.stat().st_size,
    ))
    db.commit()

    refreshed = asyncio.run(refresh_case_preparation(session, project, db))

    assert refreshed.status == "waiting_approval"
    assert refreshed.current_step == 2
    assert refreshed.artifacts["parse"]["count"] == 1
    card = refreshed.artifacts["case_prepare"]["case_cards"][0]
    assert card["readiness"] == "ready"
    assert card["modality"] == "CT"


def test_update_slide_artifacts_keeps_export_sources_in_sync():
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(engine)
    db = sessionmaker(bind=engine)()
    user = User(email="editor@example.com", name="Editor", password_hash="x")
    db.add(user)
    db.flush()
    project = Project(owner_id=user.id, title="逐页编辑", topic="验证保存", slide_count=5)
    db.add(project)
    db.commit()
    harness = AgentHarness(db, FakeServices())
    session = asyncio.run(harness.start(project, {"audience": "管理者"}))
    session = asyncio.run(harness.approve(session, project, True, ""))

    payload = SlideUpdateRequest(
        title="编辑后的第二页",
        layout="two_column",
        bullets=["第一条", "第二条"],
        key_message="编辑后的核心结论",
        speaker_notes="编辑后的讲稿备注",
    )

    apply_slide_update(session, 2, payload)

    outline_slide = session.artifacts["outline"]["slides"][1]
    generated_slide = session.artifacts["slides"]["slides"][1]
    note_item = session.artifacts["notes"]["items"][1]
    assert outline_slide["title"] == "编辑后的第二页"
    assert outline_slide["layout"] == "two_column"
    assert outline_slide["bullets"] == ["第一条", "第二条"]
    assert outline_slide["key_message"] == "编辑后的核心结论"
    assert outline_slide["speaker_notes"] == "编辑后的讲稿备注"
    assert generated_slide == outline_slide
    assert note_item == {"number": 2, "text": "编辑后的讲稿备注"}


def test_assign_slide_image_records_user_choice_for_export():
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(engine)
    db = sessionmaker(bind=engine)()
    user = User(email="image-editor@example.com", name="Image Editor", password_hash="x")
    db.add(user)
    db.flush()
    project = Project(owner_id=user.id, title="图片分配", topic="验证图片", slide_count=3)
    db.add(project)
    db.flush()
    asset = ProjectAsset(
        project_id=project.id,
        filename="chart.png",
        content_type="image/png",
        path="backend/data/uploads/demo/chart.png",
        size_bytes=1200,
    )
    db.add(asset)
    db.commit()
    harness = AgentHarness(db, FakeServices())
    session = asyncio.run(harness.start(project, {"audience": "管理者"}))
    session = asyncio.run(harness.approve(session, project, True, ""))

    apply_slide_image_assignment(session, 2, {"mode": "none"})
    assert session.artifacts["outline"]["slides"][1]["image_assignment"] == {"mode": "none"}

    apply_slide_image_assignment(session, 2, {
        "mode": "upload",
        "asset_id": asset.id,
        "filename": asset.filename,
        "path": asset.path,
        "content_type": asset.content_type,
        "author": "用户上传",
    })
    assignment = session.artifacts["outline"]["slides"][1]["image_assignment"]
    assert assignment["mode"] == "upload"
    assert assignment["asset_id"] == asset.id
    assert assignment["path"] == "backend/data/uploads/demo/chart.png"
    assert session.artifacts["slides"]["slides"][1]["image_assignment"] == assignment
    assert "_export_revision" in session.artifacts


def test_agent_services_promotes_uploaded_table_to_data_slide():
    project = Project(owner_id="owner", title="数据报告", topic="收入趋势", slide_count=5)
    artifacts = {
        "parse": {
            "items": [{
                "filename": "revenue.xlsx",
                "tables": [{
                    "sheet": "收入",
                    "rows": [["年份", "收入"], ["2024", "80"], ["2025", "100"]],
                    "data_profile": {
                        "chart_type": "line",
                        "category_column": "年份",
                        "value_columns": ["收入"],
                        "series": [{"name": "收入", "points": [{"label": "2024", "value": 80.0}, {"label": "2025", "value": 100.0}]}],
                        "kpis": [{"label": "收入", "value": 100.0, "source": "收入"}],
                        "headers": ["年份", "收入"],
                        "preview_rows": [["2024", "80"], ["2025", "100"]],
                        "source": "收入",
                    },
                }],
            }]
        }
    }

    outline = AgentServices.with_data_slides(project, AgentServices(None)._fallback_outline(project), artifacts)
    data_slide = next(slide for slide in outline["slides"] if slide.get("layout") == "data")
    assert data_slide["chart_type"] == "line"
    assert data_slide["data_source"] == {"filename": "revenue.xlsx", "sheet": "收入"}
    assert data_slide["data_table"]["headers"] == ["年份", "收入"]
    assert data_slide["data_series"][0]["name"] == "收入"


def test_radiology_outline_does_not_promote_case_index_csv_to_data_slide():
    project = Project(owner_id="owner", title="肺结节 CT 鉴别诊断病例教学", topic="病例驱动教学", slide_count=10)
    artifacts = {
        "case_prepare": {
            "case_cards": [{
                "case_id": "Case 01",
                "source_filename": "Case-01_case-note.md",
                "modality": "CT",
                "clinical_background": "65岁男性",
                "diagnosis_basis": "病理确认",
                "teaching_role": "typical",
                "display_mode": "reveal-answer-later",
                "missing_information": [],
            }],
            "case_collection_checklist": [],
            "readiness_label": "可带警告生成",
            "case_count": 1,
        },
        "parse": {
            "items": [{
                "filename": "case-index.csv",
                "tables": [{
                    "sheet": "case-index",
                    "data_profile": {
                        "chart_type": "bar",
                        "value_columns": ["key_images"],
                        "series": [{"name": "key_images", "points": [{"label": "Case-01", "value": 2.0}]}],
                    },
                }],
            }],
        },
    }

    outline = AgentServices(None)._fallback_radiology_outline(project, artifacts)

    assert not any(slide.get("layout") == "data" for slide in outline["slides"])
    assert any((slide.get("case_card") or {}).get("case_id") == "Case 01" for slide in outline["slides"])


def test_radiology_outline_follows_reference_teaching_structure():
    project = Project(owner_id="owner", title="肺结节 CT 鉴别诊断病例教学", topic="病例驱动教学", slide_count=12)
    artifacts = {
        "case_prepare": {
            "case_cards": [
                {
                    "case_id": "Case-01",
                    "modality": "CT",
                    "clinical_background": "65岁男性",
                    "diagnosis_basis": "病理确认",
                    "teaching_role": "典型病例",
                    "missing_information": [],
                },
                {
                    "case_id": "Case-02",
                    "modality": "CT",
                    "clinical_background": "随访复查",
                    "diagnosis_basis": "随访吸收",
                    "teaching_role": "鉴别诊断",
                    "missing_information": [],
                },
            ],
        },
    }

    outline = AgentServices(None)._fallback_radiology_outline(project, artifacts)
    titles = [slide["title"] for slide in outline["slides"]]

    assert titles[1:5] == ["正常解剖与检查方法基础", "疾病谱与分类框架", "影像诊断框架与读片路径", "关键征象分解"]
    assert titles.index("Case-01：典型病例读片问题") < titles.index("Case-01：影像征象与诊断推理")
    assert titles.index("Case-01：典型病例读片问题") > titles.index("关键征象分解")
    assert "跨病例对照与鉴别诊断" in titles
    assert outline["narrative"] == "正常基础—疾病谱分类—征象框架—病例推理—鉴别总结"


def test_local_case_images_match_case_slide_before_unsplash():
    artifacts = {
        "parse": {
            "items": [
                {
                    "filename": "Case-01_CT_lung-window_key-slice.png",
                    "description": "影像病例材料：radiology-case-demo/Case-01/Case-01_CT_lung-window_key-slice.png",
                    "content_type": "image/png",
                    "path": "/tmp/Case-01_CT_lung-window_key-slice.png",
                },
                {
                    "filename": "Case-02_CT_initial_nodule.png",
                    "description": "影像病例材料：radiology-case-demo/Case-02/Case-02_CT_initial_nodule.png",
                    "content_type": "image/png",
                    "path": "/tmp/Case-02_CT_initial_nodule.png",
                },
            ],
        },
    }
    slide = {"title": "典型病例：Case-01 右肺上叶结节", "layout": "content"}

    image = _local_case_image_for_slide(slide, _local_case_images(artifacts))

    assert image["filename"] == "Case-01_CT_lung-window_key-slice.png"


def test_fallback_outline_adds_architecture_for_technical_framework():
    project = Project(owner_id="owner", title="RAG Agent 技术架构", topic="讲解检索增强生成系统框架", slide_count=6)

    outline = AgentServices(None)._fallback_outline(project)

    architecture_slide = next(slide for slide in outline["slides"] if slide.get("layout") == "architecture")
    assert architecture_slide["type"] == "content"
    assert architecture_slide["diagram_title"] == "RAG Agent 技术架构"
    assert len(architecture_slide["diagram_nodes"]) >= 4
    assert len(architecture_slide["diagram_edges"]) >= 3
    assert architecture_slide["visual_rationale"]


def test_consulting_template_limits_slide_density_and_marks_template():
    outline = {
        "slides": [{
            "number": 2,
            "title": "信息过载页",
            "layout": "content",
            "bullets": ["一", "二", "三", "四", "五", "六"],
        }],
        "target_count": 1,
    }

    templated = apply_consulting_template(outline)

    assert templated["design_template"] == "consulting-default"
    assert templated["template_version"] == "1.2"
    slide = templated["slides"][0]
    assert slide["design_role"] == "content"
    assert slide["max_bullets"] == 4
    assert slide["bullets"] == ["一", "二", "三", "四"]


def test_architecture_diagram_normalizes_to_main_chain_modules():
    slide = {
        "layout": "architecture",
        "diagram_layers": ["输入", "处理", "校验", "输出", "反馈", "扩展"],
        "diagram_nodes": [
            {"id": "n1", "label": "输入", "detail": "目标与素材", "layer": "输入"},
            {"id": "n2", "label": "解析", "detail": "提取事实", "layer": "处理"},
            {"id": "n3", "label": "推理", "detail": "组织叙事", "layer": "处理"},
            {"id": "n4", "label": "校验", "detail": "来源检查", "layer": "校验"},
            {"id": "n5", "label": "生成", "detail": "页面与讲稿", "layer": "输出"},
            {"id": "n6", "label": "审批", "detail": "人工确认", "layer": "反馈"},
            {"id": "n7", "label": "归档", "detail": "保存资产", "layer": "扩展"},
        ],
    }

    normalized = normalize_architecture_diagram(slide)

    assert normalized["diagram_style"] == "main_chain"
    assert len(normalized["diagram_modules"]) <= 5
    assert [module["label"] for module in normalized["diagram_modules"]][:3] == ["输入", "处理", "校验"]
    assert all(len(module["children"]) <= 2 for module in normalized["diagram_modules"])
    assert len(normalized["diagram_edges"]) == len(normalized["diagram_modules"]) - 1


def test_export_warnings_do_not_require_images_for_architecture_slides():
    artifacts = {
        "outline": {
            "slides": [
                {"number": 1, "title": "封面", "layout": "cover", "citation_indices": [], "image_assignment": {"mode": "none"}},
                {"number": 2, "title": "技术架构", "layout": "architecture", "citation_indices": []},
            ],
        },
        "verify": {"citations": [{"title": "Source", "url": "https://example.com"}]},
        "images": {"items": []},
        "parse": {},
    }

    warnings = build_export_warnings(artifacts)

    assert not any(warning["slide_number"] == 2 and warning["code"] == "missing_image" for warning in warnings)


def test_export_manifest_exposes_template_metadata(tmp_path, monkeypatch):
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(engine)
    db = sessionmaker(bind=engine)()
    user = User(email="template@example.com", name="Template", password_hash="x")
    db.add(user)
    db.flush()
    project = Project(owner_id=user.id, title="模板验证", topic="咨询风模板", slide_count=3)
    db.add(project)
    db.commit()
    harness = AgentHarness(db, FakeServices())
    session = asyncio.run(harness.start(project, {"audience": "管理者"}))
    session = asyncio.run(harness.approve(session, project, True, ""))
    session = asyncio.run(harness.approve(session, project, True, ""))
    monkeypatch.chdir(tmp_path)
    (tmp_path / "backend" / "generated").mkdir(parents=True)

    manifest = public_export_manifest(project, session)

    assert manifest["design_template"] == "consulting-default"
    assert manifest["template_version"] == "1.2"
    assert manifest["renderer_version"]


def test_node_ppt_runtime_uses_cjk_font_for_chinese_previews():
    script = Path("scripts/ppt-runtime/generate-deck.mjs").read_text(encoding="utf-8")

    assert 'const FONT = "Heiti SC";' in script
    assert 'const FONT = "Arial";' not in script


def test_pillow_preview_renderer_can_preserve_node_layout_json(tmp_path):
    qa_dir = tmp_path / "qa"
    qa_dir.mkdir()
    layout = qa_dir / "slide-01.layout.json"
    layout.write_text('{"from":"node"}', encoding="utf-8")
    payload = {
        "slides": [{
            "number": 1,
            "title": "AI 如何改变知识工作",
            "layout": "content",
            "bullets": ["自动整理资料", "辅助形成判断"],
        }],
    }

    _render_fallback_qa(payload, qa_dir, write_layout=False)

    assert (qa_dir / "slide-01.png").exists()
    assert (qa_dir / "montage.webp").exists()
    assert layout.read_text(encoding="utf-8") == '{"from":"node"}'


def test_pillow_preview_renderer_draws_two_column_and_process_content(tmp_path):
    qa_dir = tmp_path / "qa"
    payload = {
        "slides": [
            {
                "number": 1,
                "title": "AI赋能知识工作的现状",
                "layout": "two_column",
                "left_title": "量化影响",
                "left_bullets": ["任务数量增加12.2%", "速度提高25.1%"],
                "right_title": "工作变化",
                "right_bullets": ["渗透到高技能工作", "人机协作成为新常态"],
            },
            {
                "number": 2,
                "title": "AI如何重塑工作流",
                "layout": "process",
                "process_steps": ["设定目标", "AI分析", "人工审核", "协同交付"],
            },
        ],
    }

    _render_fallback_qa(payload, qa_dir)

    from PIL import Image
    for name in ["slide-01.png", "slide-02.png"]:
        image = Image.open(qa_dir / name).convert("RGB")
        body = image.crop((80, 220, 1180, 620))
        non_white = sum(1 for pixel in body.getdata() if pixel != (255, 255, 255))
        assert non_white > 8000


def test_python_fallback_pptx_draws_two_column_and_process_content(tmp_path):
    output = tmp_path / "deck.pptx"
    qa_dir = tmp_path / "qa"
    payload = {
        "project": {"title": "模板验证", "topic": "导出正文", "slide_count": 2, "speaker_notes_enabled": True},
        "slides": [
            {
                "number": 1,
                "title": "AI赋能知识工作的现状",
                "layout": "two_column",
                "left_title": "量化影响",
                "left_bullets": ["任务数量增加12.2%", "速度提高25.1%"],
                "right_title": "工作变化",
                "right_bullets": ["渗透到高技能工作", "人机协作成为新常态"],
            },
            {
                "number": 2,
                "title": "AI如何重塑工作流",
                "layout": "process",
                "process_steps": ["设定目标", "AI分析", "人工审核", "协同交付"],
            },
        ],
    }

    _export_with_python_fallback(payload, output, qa_dir)

    from pptx import Presentation
    deck = Presentation(output)
    all_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    assert "任务数量增加12.2%" in all_text
    assert "人机协作成为新常态" in all_text
    assert "人工审核" in all_text


def test_python_fallback_pptx_preserves_speaker_notes(tmp_path):
    output = tmp_path / "deck.pptx"
    qa_dir = tmp_path / "qa"
    payload = {
        "project": {"title": "讲稿验证", "topic": "备注区", "slide_count": 1, "speaker_notes_enabled": True},
        "slides": [{
            "number": 1,
            "title": "含讲稿页面",
            "layout": "content",
            "bullets": ["正文要点"],
            "speaker_notes": "这里是需要出现在 PowerPoint 备注区的演讲稿。",
        }],
    }

    _export_with_python_fallback(payload, output, qa_dir)

    from pptx import Presentation
    deck = Presentation(output)
    notes = deck.slides[0].notes_slide.notes_text_frame.text
    assert "这里是需要出现在 PowerPoint 备注区的演讲稿。" in notes


def test_pillow_preview_font_prefers_available_cjk_font():
    source = Path("backend/app/services/pptx_export.py").read_text(encoding="utf-8")

    assert '"/System/Library/Fonts/STHeiti Medium.ttc"' in source
    assert source.index('"/System/Library/Fonts/STHeiti Medium.ttc"') < source.index('"/System/Library/Fonts/Supplemental/Arial.ttf"')


def test_export_warnings_report_missing_citations_images_and_ocr():
    artifacts = {
        "outline": {
            "slides": [
                {"number": 1, "title": "封面", "layout": "cover", "citation_indices": [], "image_assignment": {"mode": "none"}},
                {"number": 2, "title": "证据页", "layout": "evidence", "citation_indices": []},
            ],
        },
        "verify": {"citations": []},
        "images": {"items": []},
        "parse": {"ocr_required": 1, "ocr_items": ["scan.pdf"]},
    }

    warnings = build_export_warnings(artifacts)
    codes = [warning["code"] for warning in warnings]
    assert "missing_citations" in codes
    assert "missing_image" in codes
    assert "ocr_not_executed" in codes
    assert any(warning["slide_number"] == 1 and warning["code"] == "missing_image" for warning in warnings)


def test_verify_classifies_source_quality_and_deduplicates():
    research = {
        "citations": [
            {
                "title": "Peer Reviewed Study",
                "authors": ["A. Author"],
                "year": 2024,
                "venue": "Journal of Testing",
                "url": "https://doi.org/10.1000/test",
                "doi": "10.1000/test",
            },
            {
                "title": "Peer Reviewed Study",
                "authors": ["A. Author"],
                "year": 2024,
                "venue": "Journal of Testing",
                "url": "https://doi.org/10.1000/test",
                "doi": "10.1000/test",
            },
        ],
        "web_results": [
            {"title": "Industry AI Report", "url": "https://www.mckinsey.com/report", "content": "annual industry report"},
            {"title": "Breaking AI News", "url": "https://www.reuters.com/technology/ai", "content": "news update"},
            {"title": "Anonymous claim", "url": "", "content": "unattributed claim"},
        ],
    }

    verified = asyncio.run(AgentServices(None).verify(research))

    assert verified["duplicates_removed"] == 1
    assert verified["quality_summary"]["academic_paper"] == 1
    assert verified["quality_summary"]["industry_report"] == 1
    assert verified["quality_summary"]["news"] == 1
    assert verified["quality_summary"]["unreliable"] == 1
    citation = verified["citations"][0]
    assert citation["source_type"] == "academic_paper"
    assert citation["quality_tier"] == "high"
    assert citation["doi"] == "10.1000/test"
    assert citation["year"] == 2024
    assert citation["authors"] == ["A. Author"]
    assert citation["access_url"] == "https://doi.org/10.1000/test"
    assert verified["web_sources"][2]["quality_tier"] == "low"
    assert verified["web_sources"][2]["quality_notes"]


def test_literature_dedupes_across_providers_by_doi_and_title():
    citations = [
        Citation("Same DOI", ["A"], 2024, "Journal", "https://doi.org/10.1/test", "10.1/test", True, "Crossref"),
        Citation("Different title same DOI", ["B"], 2024, "Journal", "https://example.com", "10.1/TEST", True, "OpenAlex"),
        Citation("Title Only Paper", [], 2023, "Venue", "https://example.com/a", None, False, "PubMed"),
        Citation("Title-only paper!", [], 2023, "Venue", "https://example.com/b", None, False, "Semantic Scholar"),
    ]

    unique = dedupe_citations(citations)

    assert [item.provider for item in unique] == ["Crossref", "PubMed"]


def test_pubmed_helpers_extract_doi_and_year():
    doi = _pubmed_doi([
        {"idtype": "pubmed", "value": "123"},
        {"idtype": "doi", "value": "10.1000/pubmed"},
    ])

    assert doi == "10.1000/pubmed"
    assert _year_from_pubdate("2025 Jan-Feb") == 2025
