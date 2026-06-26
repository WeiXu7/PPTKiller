from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfWriter

from backend.app.services.parsers import AssetParser


def test_parses_text_and_image(tmp_path):
    parser = AssetParser()
    text_path = tmp_path / "brief.txt"
    text_path.write_text("项目背景\n核心结论", encoding="utf-8")
    text = parser.parse(str(text_path), text_path.name, "text/plain")
    assert text["parser"] == "text"
    assert "核心结论" in text["text"]

    image_path = tmp_path / "chart.png"
    Image.new("RGB", (320, 180), "white").save(image_path)
    image = parser.parse(str(image_path), image_path.name, "image/png")
    assert image["stats"]["width"] == 320
    assert image["ocr_required"] is True


def test_parses_docx_pptx_xlsx_and_pdf(tmp_path):
    parser = AssetParser()

    docx_path = tmp_path / "report.docx"
    document = Document()
    document.add_heading("行业报告", 1)
    document.add_paragraph("核心观点")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "指标"
    table.cell(0, 1).text = "数值"
    table.cell(1, 0).text = "增长率"
    table.cell(1, 1).text = "20%"
    document.save(docx_path)
    docx = parser.parse(str(docx_path), docx_path.name, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    assert docx["stats"]["paragraphs"] == 2
    assert docx["stats"]["tables"] == 1
    assert "行业报告" in docx["text"]

    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "市场机会"
    slide.placeholders[1].text = "机会一\n机会二"
    presentation.save(pptx_path)
    pptx = parser.parse(str(pptx_path), pptx_path.name, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    assert pptx["stats"]["slides"] == 1
    assert "市场机会" in pptx["text"]

    xlsx_path = tmp_path / "data.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "数据"
    sheet.append(["年份", "收入"])
    sheet.append([2025, 100])
    workbook.save(xlsx_path)
    xlsx = parser.parse(str(xlsx_path), xlsx_path.name, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    assert xlsx["stats"]["sheets"] == 1
    assert "收入" in xlsx["text"]
    profile = xlsx["tables"][0]["data_profile"]
    assert profile["chart_type"] == "bar"
    assert profile["category_column"] == "年份"
    assert profile["value_columns"] == ["收入"]
    assert profile["series"][0]["points"] == [{"label": "2025", "value": 100.0}]

    pdf_path = tmp_path / "blank.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with pdf_path.open("wb") as output:
        writer.write(output)
    pdf = parser.parse(str(pdf_path), pdf_path.name, "application/pdf")
    assert pdf["stats"]["pages"] == 1
